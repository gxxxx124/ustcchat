#!/usr/bin/env python3
"""
轻量级marker转换器，避免下载大模型
"""

import os
import logging
from typing import Dict, Any, Optional
from pathlib import Path

# 配置日志
logger = logging.getLogger("lightweight_marker")
logger.setLevel(logging.INFO)
if not logger.handlers:
    handler = logging.StreamHandler()
    formatter = logging.Formatter('%(asctime)s - %(name)s - %(levelname)s - %(message)s')
    handler.setFormatter(formatter)
    logger.addHandler(handler)

def convert_with_lightweight_marker(
    file_path: str, 
    output_dir: Optional[str] = None,
    base_name: Optional[str] = None
) -> Dict[str, Any]:
    """
    使用轻量级marker转换文档
    """
    try:
        logger.info(f"🔄 开始轻量级转换: {file_path}")
        
        # 检查文件是否存在
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"文件不存在: {file_path}")
        
        # 设置输出目录
        if output_dir is None:
            output_dir = os.path.dirname(os.path.abspath(file_path))
        
        # 创建输出目录
        os.makedirs(output_dir, exist_ok=True)
        
        # 生成输出文件名
        if base_name is None:
            base_name = os.path.splitext(os.path.basename(file_path))[0]
        
        # 根据文件类型选择处理方法
        file_ext = os.path.splitext(file_path)[1].lower()
        
        if file_ext == '.pdf':
            return convert_pdf_lightweight(file_path, output_dir, base_name)
        elif file_ext in ['.docx', '.doc']:
            return convert_word_lightweight(file_path, output_dir, base_name)
        elif file_ext in ['.pptx', '.ppt']:
            return convert_powerpoint_lightweight(file_path, output_dir, base_name)
        elif file_ext in ['.xlsx', '.xls']:
            return convert_excel_lightweight(file_path, output_dir, base_name)
        else:
            raise ValueError(f"不支持的文件类型: {file_ext}")
            
    except Exception as e:
        logger.error(f"❌ 轻量级转换失败: {str(e)}")
        return {
            "success": False,
            "message": f"转换失败: {str(e)}",
            "data": {
                "file_path": file_path,
                "error": str(e)
            }
        }

def convert_pdf_lightweight(file_path: str, output_dir: str, base_name: str) -> Dict[str, Any]:
    """PDF轻量级转换"""
    try:
        import pdfplumber
        
        text_content = []
        with pdfplumber.open(file_path) as pdf:
            total_pages = len(pdf.pages)
            logger.info(f"📄 PDF共有 {total_pages} 页")
            
            for page_num, page in enumerate(pdf.pages, 1):
                # 提取文本
                page_text = page.extract_text()
                if page_text and page_text.strip():
                    text_content.append(f"## 第 {page_num} 页\n\n{page_text.strip()}\n")
                
                # 提取表格
                tables = page.extract_tables()
                if tables:
                    for table_num, table in enumerate(tables, 1):
                        if table and len(table) > 1:
                            table_md = f"\n### 表格 {table_num}\n\n"
                            # 转换表格为Markdown
                            if len(table) > 0:
                                header = table[0]
                                if header:
                                    table_md += "| " + " | ".join(str(cell or "") for cell in header) + " |\n"
                                    table_md += "| " + " | ".join("---" for _ in header) + " |\n"
                                
                                for row in table[1:]:
                                    if row:
                                        table_md += "| " + " | ".join(str(cell or "") for cell in row) + " |\n"
                            
                            text_content.append(table_md + "\n")
        
        markdown_content = "\n".join(text_content)
        
        # 保存文件
        md_file_path = os.path.join(output_dir, f"{base_name}.md")
        with open(md_file_path, 'w', encoding='utf-8') as f:
            f.write(markdown_content)
        
        return {
            "success": True,
            "message": f"PDF转换成功: {file_path}",
            "data": {
                "file_path": file_path,
                "output_dir": output_dir,
                "text_length": len(markdown_content),
                "method": "pdfplumber"
            }
        }
        
    except Exception as e:
        raise Exception(f"PDF转换失败: {str(e)}")

def convert_word_lightweight(file_path: str, output_dir: str, base_name: str) -> Dict[str, Any]:
    """Word文档轻量级转换"""
    try:
        from docx import Document
        
        doc = Document(file_path)
        text_content = []
        
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_content.append(paragraph.text.strip())
        
        # 处理表格
        for table in doc.tables:
            table_md = "\n### 表格\n\n"
            for row in table.rows:
                row_text = "| " + " | ".join(cell.text.strip() for cell in row.cells) + " |"
                table_md += row_text + "\n"
            text_content.append(table_md)
        
        markdown_content = "\n\n".join(text_content)
        
        # 保存文件
        md_file_path = os.path.join(output_dir, f"{base_name}.md")
        with open(md_file_path, 'w', encoding='utf-8') as f:
            f.write(markdown_content)
        
        return {
            "success": True,
            "message": f"Word转换成功: {file_path}",
            "data": {
                "file_path": file_path,
                "output_dir": output_dir,
                "text_length": len(markdown_content),
                "method": "python-docx"
            }
        }
        
    except Exception as e:
        raise Exception(f"Word转换失败: {str(e)}")

def convert_powerpoint_lightweight(file_path: str, output_dir: str, base_name: str) -> Dict[str, Any]:
    """PowerPoint轻量级转换"""
    try:
        from pptx import Presentation
        
        prs = Presentation(file_path)
        text_content = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = f"## 幻灯片 {slide_num}\n\n"
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text += shape.text.strip() + "\n"
            
            if slide_text.strip():
                text_content.append(slide_text)
        
        markdown_content = "\n".join(text_content)
        
        # 保存文件
        md_file_path = os.path.join(output_dir, f"{base_name}.md")
        with open(md_file_path, 'w', encoding='utf-8') as f:
            f.write(markdown_content)
        
        return {
            "success": True,
            "message": f"PowerPoint转换成功: {file_path}",
            "data": {
                "file_path": file_path,
                "output_dir": output_dir,
                "text_length": len(markdown_content),
                "method": "python-pptx"
            }
        }
        
    except Exception as e:
        raise Exception(f"PowerPoint转换失败: {str(e)}")

def convert_excel_lightweight(file_path: str, output_dir: str, base_name: str) -> Dict[str, Any]:
    """Excel轻量级转换"""
    try:
        import pandas as pd
        
        # 读取所有工作表
        excel_file = pd.ExcelFile(file_path)
        text_content = []
        
        for sheet_name in excel_file.sheet_names:
            df = pd.read_excel(file_path, sheet_name=sheet_name)
            
            if not df.empty:
                text_content.append(f"## 工作表: {sheet_name}\n")
                
                # 转换为Markdown表格
                markdown_table = df.to_markdown(index=False)
                text_content.append(markdown_table)
                text_content.append("\n")
        
        markdown_content = "\n".join(text_content)
        
        # 保存文件
        md_file_path = os.path.join(output_dir, f"{base_name}.md")
        with open(md_file_path, 'w', encoding='utf-8') as f:
            f.write(markdown_content)
        
        return {
            "success": True,
            "message": f"Excel转换成功: {file_path}",
            "data": {
                "file_path": file_path,
                "output_dir": output_dir,
                "text_length": len(markdown_content),
                "method": "pandas"
            }
        }
        
    except Exception as e:
        raise Exception(f"Excel转换失败: {str(e)}")

if __name__ == "__main__":
    import sys
    
    if len(sys.argv) < 2:
        print("用法: python lightweight_marker_converter.py <文件路径> [输出目录]")
        sys.exit(1)
    
    file_path = sys.argv[1]
    output_dir = sys.argv[2] if len(sys.argv) > 2 else None
    
    result = convert_with_lightweight_marker(file_path, output_dir)
    print(f"转换结果: {result}")
